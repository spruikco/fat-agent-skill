#!/usr/bin/env python3
"""
generate-report.py — Generate Word (.docx) and PowerPoint (.pptx) audit reports.

Reads scored JSON (from calculate-score.py), optional SEMrush data, and
optional chart images to produce professional branded reports.

Usage:
    # Generate both reports from scored JSON
    python analyse-html.py page.html | python calculate-score.py | \\
        python generate-report.py --url example.com --output-dir ./reports

    # With SEMrush data and charts
    python generate-report.py --scores scores.json --semrush semrush.json \\
        --charts-dir ./charts --brand logo.png --output-dir ./reports

    # Word only
    python generate-report.py --scores scores.json --format docx

Options:
    --scores FILE       Path to scored JSON (or pipe via stdin)
    --semrush FILE      Path to supplementary SEMrush data JSON
    --url URL           The audited site URL (for report title)
    --output-dir DIR    Output directory (default: ./reports)
    --format FORMAT     Output format: docx, pptx, html, or both (default: both)
    --charts-dir DIR    Directory containing chart PNG images
    --brand IMAGE       Path to brand/logo image for cover pages
    --font FONT         Font family name (default: Plus Jakarta Sans)
    --pagespeed FILE    Path to PageSpeed data JSON (mobile + desktop)
    --client-facing     Rewrite findings into plain-English, jargon-free language

Dependencies:
    python-docx (pip install python-docx)
    python-pptx (pip install python-pptx)
    Pillow      (pip install Pillow)     — optional, for image handling

Output:
    FAT_Audit_Report.docx       — Comprehensive Word report
    FAT_Audit_Presentation.pptx — Executive summary PowerPoint
"""

import sys
import json
import os
import re
import argparse
import subprocess
from datetime import date

# Make sibling scripts importable (client_facing helpers).
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))


def score_to_grade(score):
    """Single source of truth for letter grades.

    Bands match calculate-score.py: A>=90, B>=75, C>=60, D>=40, else F.
    Used by BOTH the Word score table and the PPTX score cards so a given
    score always maps to the same grade across deliverables.
    """
    try:
        s = float(score)
    except (TypeError, ValueError):
        s = 0
    if s >= 90:
        return "A"
    if s >= 75:
        return "B"
    if s >= 60:
        return "C"
    if s >= 40:
        return "D"
    return "F"


def _category_score(scores, category):
    """Read a category score from the nested schema, falling back to the
    flat `<category>_score` key (sample_scores.json shape). Prefers nested."""
    val = scores.get(category)
    if isinstance(val, dict) and val.get("score") is not None:
        return val.get("score")
    return scores.get(f"{category}_score", 0) or 0


def safe_report_slug(url):
    """Build a filesystem-safe slug from an arbitrary URL.

    Strips path separators and traversal sequences so a malicious --url
    cannot escape the output directory. Returns '' for falsy input.
    """
    if not url:
        return ""
    slug = re.sub(r"[^A-Za-z0-9._-]", "_", url)
    # basename collapses any residual separators; guards against '..' too.
    slug = os.path.basename(slug)
    # A bare '..' (or empty) is not a usable component.
    if slug in ("", ".", ".."):
        return ""
    return slug


def _safe_output_path(output_dir, filename):
    """Join output_dir/filename and assert the realpath stays inside output_dir."""
    base = os.path.realpath(output_dir)
    candidate = os.path.realpath(os.path.join(base, filename))
    if candidate != base and not candidate.startswith(base + os.sep):
        raise ValueError(
            f"Refusing to write outside output directory: {candidate!r} not under {base!r}"
        )
    return candidate


def _client_facing_string(text):
    """Apply client-facing transforms to a single finding string.

    Reuses scripts/client_facing.py helpers: strip code snippets, then relabel
    jargon (incl. P0-P3 -> Urgent/Important/Recommended/Nice to Have).
    """
    from client_facing import transform_text, strip_code_blocks

    return transform_text(strip_code_blocks(str(text))).strip()


def make_client_facing(scores):
    """Return a copy of the scored data with findings/labels made client-friendly.

    Transforms the canonical `summary.{critical,high,medium,low}` string lists
    AND the flat `findings[]` objects (sample_scores.json shape), stripping code
    snippets and replacing technical jargon BEFORE the data is rendered into any
    deliverable (docx, pptx, html).
    """
    import copy

    from client_facing import transform_text, strip_code_blocks

    data = copy.deepcopy(scores)

    summary = data.get("summary")
    if isinstance(summary, dict):
        for bucket in ("critical", "high", "medium", "low"):
            items = summary.get(bucket)
            if isinstance(items, list):
                summary[bucket] = [_client_facing_string(it) for it in items]

    flat = data.get("findings")
    if isinstance(flat, list):
        for f in flat:
            if not isinstance(f, dict):
                continue
            if "priority" in f:
                f["priority"] = transform_text(str(f["priority"]))
            for key in ("issue", "title", "category"):
                if key in f and isinstance(f[key], str):
                    f[key] = transform_text(strip_code_blocks(f[key])).strip()
            for key in ("detail", "fix"):
                if key in f and isinstance(f[key], str):
                    f[key] = transform_text(strip_code_blocks(f[key])).strip()

    return data


# Defer imports for helpful error messages
MISSING_DEPS = []
try:
    from docx import Document
    from docx.shared import Inches, Pt, Cm, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    from docx.oxml.ns import nsdecls
    from docx.oxml import parse_xml
except ImportError:
    MISSING_DEPS.append("python-docx")

try:
    from pptx import Presentation
    from pptx.util import Inches as PInches, Pt as PPt
    from pptx.dml.color import RGBColor as PRGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    MISSING_DEPS.append("python-pptx")

# ---------- Constants ----------
FONT = "Plus Jakarta Sans"

# Word colors
DARK_GRAY = RGBColor(0x2C, 0x3E, 0x50) if "python-docx" not in MISSING_DEPS else None
ACCENT_RED = RGBColor(0xC0, 0x39, 0x2B) if "python-docx" not in MISSING_DEPS else None
MID_GRAY = RGBColor(0x7F, 0x8C, 0x8D) if "python-docx" not in MISSING_DEPS else None
WHITE = RGBColor(0xFF, 0xFF, 0xFF) if "python-docx" not in MISSING_DEPS else None

# PPTX colors
if "python-pptx" not in MISSING_DEPS:
    P_DARK = PRGBColor(0x1A, 0x1A, 0x2E)
    P_RED = PRGBColor(0xC0, 0x39, 0x2B)
    P_ORANGE = PRGBColor(0xE6, 0x7E, 0x22)
    P_YELLOW = PRGBColor(0xF3, 0x9C, 0x12)
    P_GREEN = PRGBColor(0x27, 0xAE, 0x60)
    P_WHITE = PRGBColor(0xFF, 0xFF, 0xFF)
    P_GRAY = PRGBColor(0x2C, 0x3E, 0x50)
    P_MGRAY = PRGBColor(0x7F, 0x8C, 0x8D)
    P_LGRAY = PRGBColor(0xEC, 0xF0, 0xF1)
    P_BLUE = PRGBColor(0x29, 0x80, 0xB9)


# ---------- Helper Functions ----------


def _set_cell_shading(cell, color_hex):
    shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{color_hex}"/>')
    cell._tc.get_or_add_tcPr().append(shading)


def _set_cell(cell, text, bold=False, color=None, size=10, alignment=None, font=FONT):
    cell.text = ""
    p = cell.paragraphs[0]
    if alignment:
        p.alignment = alignment
    run = p.add_run(str(text))
    run.bold = bold
    run.font.size = Pt(size)
    run.font.name = font
    if color:
        run.font.color.rgb = color
    p.paragraph_format.space_before = Pt(2)
    p.paragraph_format.space_after = Pt(2)


def _make_table(doc, headers):
    table = doc.add_table(rows=1, cols=len(headers))
    table.style = "Table Grid"
    for i, h in enumerate(headers):
        _set_cell(
            table.rows[0].cells[i],
            h,
            bold=True,
            color=WHITE,
            size=9,
            alignment=WD_ALIGN_PARAGRAPH.CENTER,
        )
        _set_cell_shading(table.rows[0].cells[i], "2C3E50")
    return table


def _add_check_row(table, check, result, detail, status):
    row = table.add_row()
    for i, val in enumerate([check, result, detail, status]):
        _set_cell(
            row.cells[i],
            val,
            size=8,
            bold=(i == 0),
            alignment=WD_ALIGN_PARAGRAPH.CENTER if i > 0 else None,
        )
    colors = {"PASS": "D5F5E3", "FAIL": "FADBD8", "WARNING": "FDEBD0", "N/A": "F2F3F4"}
    _set_cell_shading(row.cells[3], colors.get(status, "F2F3F4"))


def _heading(doc, text, level=1):
    h = doc.add_heading(text, level=level)
    for r in h.runs:
        r.font.color.rgb = DARK_GRAY
        r.font.name = FONT
    return h


def _para(doc, text, bold=False, color=None, size=10):
    p = doc.add_paragraph()
    run = p.add_run(text)
    run.bold = bold
    run.font.size = Pt(size)
    run.font.name = FONT
    if color:
        run.font.color.rgb = color
    return p


def _insert_chart(doc, charts_dir, filename, caption, width=6.0):
    """Insert a chart image with caption if the file exists."""
    if not charts_dir:
        return False
    path = os.path.join(charts_dir, filename)
    if not os.path.exists(path):
        return False
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.add_run().add_picture(path, width=Inches(width))
    cap = doc.add_paragraph()
    cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = cap.add_run(caption)
    run.font.size = Pt(9)
    run.font.italic = True
    run.font.name = FONT
    return True


# ---------- PPTX Helpers ----------


def _ptb(
    slide,
    left,
    top,
    width,
    height,
    text,
    size=18,
    color=P_WHITE,
    bold=False,
    align=PP_ALIGN.LEFT,
):
    txBox = slide.shapes.add_textbox(
        PInches(left), PInches(top), PInches(width), PInches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = PPt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = FONT
    p.alignment = align


def _prect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        PInches(left),
        PInches(top),
        PInches(width),
        PInches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _pcard(slide, left, top, width, height, label, score, color):
    """Render a score card: the numeric ``score`` is the large text and the
    category ``label`` is the small caption beneath it."""
    _prect(slide, left, top, width, height, color)
    _ptb(
        slide,
        left,
        top + 0.12,
        width,
        height * 0.55,
        str(score),
        28,
        P_WHITE,
        True,
        PP_ALIGN.CENTER,
    )
    _ptb(
        slide,
        left,
        top + height * 0.55,
        width,
        height * 0.4,
        str(label),
        10,
        P_WHITE,
        True,
        PP_ALIGN.CENTER,
    )


def _pheader(slide, title, brand=None):
    _prect(slide, 0, 0, 13.333, 1.1, P_DARK)
    if brand and os.path.exists(brand):
        slide.shapes.add_picture(
            brand, PInches(0.3), PInches(0.05), PInches(1.0), PInches(1.0)
        )
    _ptb(slide, 1.4 if brand else 0.5, 0.25, 10, 0.6, title, 28, P_WHITE, True)


def _pbullet(slide, x, y, text, size, text_color, dot_color):
    _prect(slide, x, y + 0.07, 0.1, 0.1, dot_color)
    _ptb(slide, x + 0.15, y, 11.0, 0.3, text, size, text_color)


def _add_chart_fit(slide, path, left, top, max_w, max_h):
    """Place a chart image scaled to FIT inside a (max_w x max_h) inch box while
    preserving its native aspect ratio, then centre it in that box.

    The old code forced every chart into a fixed width AND height, which
    stretched/squashed charts whose real aspect ratio differed from the box
    (the "stretched graphs" problem). Reading the true pixel dimensions and
    scaling by the limiting axis keeps every chart undistorted.
    """
    iw = ih = None
    try:
        from PIL import Image

        with Image.open(path) as im:
            iw, ih = im.size
    except Exception:
        iw = ih = None

    if iw and ih:
        ratio = iw / ih
        box_ratio = max_w / max_h
        if ratio >= box_ratio:
            w, h = max_w, max_w / ratio
        else:
            h, w = max_h, max_h * ratio
        cx = left + (max_w - w) / 2
        cy = top + (max_h - h) / 2
        return slide.shapes.add_picture(
            path, PInches(cx), PInches(cy), PInches(w), PInches(h)
        )
    # Unknown dimensions (no Pillow): constrain by width only — python-pptx then
    # derives the height from the image's own ratio, so it still won't stretch.
    return slide.shapes.add_picture(
        path, PInches(left), PInches(top), width=PInches(max_w)
    )


def _ptable(
    slide,
    left,
    top,
    width,
    headers,
    rows,
    col_widths=None,
    header_color=None,
    font_size=10,
    row_height=0.34,
):
    """Render a styled table on a slide. Explicit per-cell fills override the
    default banded table style so colours stay on-brand and legible."""
    header_color = header_color or P_DARK
    nrows, ncols = len(rows) + 1, len(headers)
    gfx = slide.shapes.add_table(
        nrows,
        ncols,
        PInches(left),
        PInches(top),
        PInches(width),
        PInches(row_height * nrows),
    )
    table = gfx.table
    if col_widths and len(col_widths) == ncols:
        total = float(sum(col_widths))
        for i, cw in enumerate(col_widths):
            table.columns[i].width = PInches(width * cw / total)

    def _fill(cell, text, color, bold=False):
        cell.fill.solid()
        cell.fill.fore_color.rgb = color
        cell.margin_top = PInches(0.02)
        cell.margin_bottom = PInches(0.02)
        para = cell.text_frame.paragraphs[0]
        para.text = str(text)
        para.font.size = PPt(font_size)
        para.font.bold = bold
        para.font.name = FONT
        para.font.color.rgb = P_WHITE if bold else P_GRAY

    for j, h in enumerate(headers):
        _fill(table.cell(0, j), h, header_color, bold=True)
    for i, row in enumerate(rows, start=1):
        band = P_WHITE if i % 2 else P_LGRAY
        for j, val in enumerate(row):
            _fill(table.cell(i, j), val, band)
    return table


def _strip_host(u):
    """Reduce a full URL to its path so target pages read cleanly in tables."""
    p = re.sub(r"^https?://[^/]+", "", str(u or ""))
    return p or "/"


# ---------- Word Report Generator ----------


def generate_docx(
    scores,
    url,
    output_dir,
    charts_dir=None,
    semrush=None,
    brand=None,
    pagespeed=None,
    font=FONT,
    recommendations=None,
):
    global FONT
    FONT = font

    doc = Document()
    style = doc.styles["Normal"]
    style.font.name = FONT
    style.font.size = Pt(10)
    style.font.color.rgb = DARK_GRAY

    for section in doc.sections:
        section.top_margin = Cm(2)
        section.bottom_margin = Cm(2)
        section.left_margin = Cm(2.5)
        section.right_margin = Cm(2.5)

    # --- Cover ---
    for _ in range(3):
        doc.add_paragraph()
    if brand and os.path.exists(brand):
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.add_run().add_picture(brand, width=Inches(3))
        doc.add_paragraph()

    for text, sz, clr, bld in [
        ("FAT AGENT AUDIT REPORT", 28, DARK_GRAY, True),
        ("Fix  |  Audit  |  Test", 16, MID_GRAY, False),
        ("", 10, None, False),
        (url or "Website Audit", 20, ACCENT_RED, True),
        (f"Audit Date: {date.today().strftime('%d %B %Y')}", 12, MID_GRAY, False),
        ("CONFIDENTIAL", 10, ACCENT_RED, True),
    ]:
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        if text:
            run = p.add_run(text)
            run.bold = bld
            run.font.size = Pt(sz)
            run.font.name = FONT
            if clr:
                run.font.color.rgb = clr
    doc.add_page_break()

    # --- Scoring Summary ---
    _heading(doc, "1. Scoring Summary")

    overall = scores.get("overall", {})
    seo_score = _category_score(scores, "seo")
    sec_score = _category_score(scores, "security")
    a11y_score = _category_score(scores, "accessibility")
    perf_score = _category_score(scores, "performance")
    fat_score = overall.get("score", scores.get("overall_score", 0)) or 0

    score_table = _make_table(doc, ["Category", "FAT Score", "Grade"])
    for cat, sc in [
        ("SEO", seo_score),
        ("Security", sec_score),
        ("Accessibility", a11y_score),
        ("Performance", perf_score),
        ("Overall", fat_score),
    ]:
        row = score_table.add_row()
        # Single shared grade function — identical bands to the PPTX cards.
        g = score_to_grade(sc)
        for i, v in enumerate([cat, f"{sc}/100", g]):
            _set_cell(
                row.cells[i],
                v,
                size=10,
                bold=(i == 0),
                alignment=WD_ALIGN_PARAGRAPH.CENTER,
            )
        grade_colors = {
            "A": "D5F5E3",
            "B": "D6EAF8",
            "C": "FDEBD0",
            "D": "F5CBA7",
            "F": "FADBD8",
        }
        _set_cell_shading(row.cells[2], grade_colors.get(g, "F2F3F4"))

    _insert_chart(
        doc,
        charts_dir,
        "chart_fat_scores.png",
        "Figure: FAT Agent Score Breakdown & Issue Distribution",
        6.0,
    )

    doc.add_page_break()

    # --- Findings ---
    _heading(doc, "2. All Findings")
    summary = scores.get("summary", {})
    all_issues = []
    for prio, label in [
        ("critical", "P0"),
        ("high", "P1"),
        ("medium", "P2"),
        ("low", "P3"),
    ]:
        for issue in summary.get(prio, []):
            all_issues.append((label, issue))

    if all_issues:
        findings_table = _make_table(doc, ["#", "Priority", "Finding"])
        for idx, (prio, issue) in enumerate(all_issues, 1):
            row = findings_table.add_row()
            _set_cell(
                row.cells[0], str(idx), size=9, alignment=WD_ALIGN_PARAGRAPH.CENTER
            )
            _set_cell(
                row.cells[1],
                prio,
                bold=True,
                size=9,
                color=WHITE,
                alignment=WD_ALIGN_PARAGRAPH.CENTER,
            )
            _set_cell(row.cells[2], issue, size=8)
            color_map = {"P0": "C0392B", "P1": "E67E22", "P2": "F39C12", "P3": "27AE60"}
            _set_cell_shading(row.cells[1], color_map.get(prio, "7F8C8D"))
    else:
        _para(
            doc,
            "No issues found — congratulations!",
            color=RGBColor(0x27, 0xAE, 0x60),
            bold=True,
        )

    doc.add_page_break()

    # --- SEO Breakdown ---
    _heading(doc, "3. SEO Score Breakdown")
    seo_details = scores.get("seo", {}).get("details", {})
    if seo_details:
        seo_table = _make_table(doc, ["Category", "Score", "Max"])
        for cat, data in seo_details.items():
            row = seo_table.add_row()
            label = cat.replace("_", " ").title()
            sc = data.get("score", 0)
            mx = data.get("max", 0)
            _set_cell(row.cells[0], label, bold=True, size=9)
            _set_cell(
                row.cells[1], str(sc), size=9, alignment=WD_ALIGN_PARAGRAPH.CENTER
            )
            _set_cell(
                row.cells[2], str(mx), size=9, alignment=WD_ALIGN_PARAGRAPH.CENTER
            )
            if sc >= mx:
                _set_cell_shading(row.cells[1], "D5F5E3")
            elif sc >= mx * 0.6:
                _set_cell_shading(row.cells[1], "FDEBD0")
            else:
                _set_cell_shading(row.cells[1], "FADBD8")

    doc.add_page_break()

    # --- SEMrush Section ---
    if semrush:
        _heading(doc, "4. SEMrush Domain Intelligence")
        domain = semrush.get("domain", url or "")
        sem_table = doc.add_table(rows=0, cols=2)
        sem_table.style = "Table Grid"
        for label, value in [
            ("Domain", domain),
            ("Authority Score", f"{semrush.get('authority_score', '?')}/100"),
            (
                "Organic Traffic",
                f"{semrush.get('organic_traffic', '?')}/month ({semrush.get('traffic_change', '')})",
            ),
            (
                "Organic Keywords",
                f"{semrush.get('organic_keywords', '?')} ({semrush.get('keywords_change', '')})",
            ),
            ("Referring Domains", str(semrush.get("referring_domains", "?"))),
            ("Backlinks", str(semrush.get("backlinks", "?"))),
        ]:
            row = sem_table.add_row()
            _set_cell(row.cells[0], label, bold=True, size=10)
            _set_cell(row.cells[1], value, size=10)
            _set_cell_shading(row.cells[0], "F2F3F4")

        doc.add_paragraph()
        _insert_chart(
            doc,
            charts_dir,
            "chart_overview.png",
            "Figure: Domain Overview Metrics",
            6.5,
        )
        _insert_chart(
            doc,
            charts_dir,
            "chart_traffic_trend.png",
            "Figure: Organic Traffic Trend",
            6.0,
        )
        _insert_chart(
            doc,
            charts_dir,
            "chart_keywords_trend.png",
            "Figure: Keywords Trend & SERP Distribution",
            6.0,
        )
        _insert_chart(
            doc,
            charts_dir,
            "chart_top_keywords.png",
            "Figure: Top Keywords by Volume",
            5.5,
        )

        # Priority opportunities table (striking-distance, high-value keywords)
        opps = semrush.get("opportunity_keywords")
        if opps:
            _para(
                doc,
                "Priority Keyword Opportunities",
                bold=True,
                size=12,
                color=DARK_GRAY,
            )
            _para(
                doc,
                "High-value keywords within striking distance, ranked by volume x CPC and winnability.",
                size=9,
                color=MID_GRAY,
            )
            opp_table = _make_table(
                doc, ["Keyword", "Vol", "CPC", "Pos", "Target Page", "Priority"]
            )
            for k in opps[:15]:
                try:
                    cpc_s = f"${float(k.get('cpc')):.2f}"
                except (TypeError, ValueError):
                    cpc_s = "-"
                r = opp_table.add_row()
                _set_cell(r.cells[0], str(k.get("keyword", "")), size=9, bold=True)
                _set_cell(
                    r.cells[1],
                    str(k.get("volume", "")),
                    size=9,
                    alignment=WD_ALIGN_PARAGRAPH.CENTER,
                )
                _set_cell(
                    r.cells[2], cpc_s, size=9, alignment=WD_ALIGN_PARAGRAPH.CENTER
                )
                _set_cell(
                    r.cells[3],
                    str(k.get("position", "")),
                    size=9,
                    alignment=WD_ALIGN_PARAGRAPH.CENTER,
                )
                _set_cell(r.cells[4], _strip_host(k.get("url", "")), size=8)
                _set_cell(
                    r.cells[5],
                    str(k.get("priority", "")),
                    size=9,
                    alignment=WD_ALIGN_PARAGRAPH.CENTER,
                )
            doc.add_paragraph()

        # Cannibalisation table
        cann = semrush.get("cannibalization")
        if cann:
            _para(doc, "Keyword Cannibalisation", bold=True, size=12, color=DARK_GRAY)
            _para(
                doc,
                "Same keyword ranking on multiple URLs — pages compete and split ranking signals.",
                size=9,
                color=MID_GRAY,
            )
            cann_table = _make_table(
                doc, ["Keyword", "Vol", "# URLs", "Competing Pages"]
            )
            for c in cann[:12]:
                urls = c.get("urls", []) or []
                count = len(urls) if urls else c.get("url_count", "")
                r = cann_table.add_row()
                _set_cell(r.cells[0], str(c.get("keyword", "")), size=9, bold=True)
                _set_cell(
                    r.cells[1],
                    str(c.get("volume", "")),
                    size=9,
                    alignment=WD_ALIGN_PARAGRAPH.CENTER,
                )
                _set_cell(
                    r.cells[2], str(count), size=9, alignment=WD_ALIGN_PARAGRAPH.CENTER
                )
                _set_cell(
                    r.cells[3], "  ".join(_strip_host(u) for u in urls[:4]), size=8
                )

        doc.add_page_break()

    # --- Recommended Action Plan ---
    if recommendations:
        _heading(doc, "Recommended Action Plan")
        if isinstance(recommendations[0], dict):
            phases = [
                (r.get("phase") or r.get("title") or "", r.get("items") or [])
                for r in recommendations
            ]
        else:
            phases = [("Recommended Actions", list(recommendations))]
        for pname, items in phases:
            if pname:
                _para(doc, pname, bold=True, size=12, color=ACCENT_RED)
            for it in items:
                p = doc.add_paragraph(style="List Bullet")
                run = p.add_run(str(it))
                run.font.size = Pt(10)
                run.font.name = FONT
                run.font.color.rgb = DARK_GRAY
        doc.add_page_break()

    # --- PageSpeed ---
    if pagespeed:
        _heading(doc, "5. PageSpeed Performance")
        _insert_chart(
            doc,
            charts_dir,
            "chart_pagespeed.png",
            "Figure: PageSpeed Mobile vs Desktop",
            5.5,
        )

    # --- Footer ---
    doc.add_paragraph()
    if brand and os.path.exists(brand):
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.add_run().add_picture(brand, width=Inches(1.2))
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(
        f"Report generated by FAT Agent | {date.today().strftime('%d %B %Y')}"
    )
    run.font.size = Pt(8)
    run.font.color.rgb = MID_GRAY
    run.font.name = FONT

    slug = safe_report_slug(url)
    filename = f"FAT_Audit_Report_{slug}.docx" if slug else "FAT_Audit_Report.docx"
    path = _safe_output_path(output_dir, filename)
    doc.save(path)
    return path


# ---------- PPTX Report Generator ----------


def generate_pptx(
    scores,
    url,
    output_dir,
    charts_dir=None,
    semrush=None,
    brand=None,
    pagespeed=None,
    font=FONT,
    recommendations=None,
):
    global FONT
    FONT = font

    prs = Presentation()
    prs.slide_width = PInches(13.333)
    prs.slide_height = PInches(7.5)

    overall = scores.get("overall", {})
    seo_score = _category_score(scores, "seo")
    sec_score = _category_score(scores, "security")
    a11y_score = _category_score(scores, "accessibility")
    perf_score = _category_score(scores, "performance")
    fat_score = overall.get("score", scores.get("overall_score", 0)) or 0

    # Colour cards by the SHARED grade so the colour matches the letter grade
    # shown for the same score in the Word report.
    _GRADE_CARD_COLOR = {
        "A": P_GREEN,
        "B": P_GREEN,
        "C": P_ORANGE,
        "D": P_ORANGE,
        "F": P_RED,
    }

    def _grade_color(sc):
        return _GRADE_CARD_COLOR.get(score_to_grade(sc), P_RED)

    def add_bg(slide, color=P_WHITE):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = color

    # --- Slide 1: Title ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, P_DARK)
    if brand and os.path.exists(brand):
        slide.shapes.add_picture(
            brand, PInches(4.9), PInches(0.3), PInches(3.5), PInches(3.5)
        )
    _ptb(
        slide,
        0.5,
        4.0,
        12.3,
        1.0,
        "FAT AGENT AUDIT REPORT",
        40,
        P_WHITE,
        True,
        PP_ALIGN.CENTER,
    )
    _ptb(
        slide,
        0.5,
        4.8,
        12.3,
        0.5,
        "Fix  |  Audit  |  Test",
        20,
        P_MGRAY,
        False,
        PP_ALIGN.CENTER,
    )
    _ptb(
        slide,
        0.5,
        5.5,
        12.3,
        0.6,
        url or "Website Audit",
        28,
        P_RED,
        True,
        PP_ALIGN.CENTER,
    )
    _ptb(
        slide,
        0.5,
        6.3,
        12.3,
        0.5,
        date.today().strftime("%d %B %Y"),
        14,
        P_MGRAY,
        False,
        PP_ALIGN.CENTER,
    )

    # --- Slide 2: Executive Summary ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    _pheader(slide, "Executive Summary", brand)
    _pcard(
        slide,
        0.3,
        1.5,
        2.4,
        1.3,
        "SEO",
        f"{seo_score}/100",
        _grade_color(seo_score),
    )
    _pcard(
        slide,
        2.9,
        1.5,
        2.4,
        1.3,
        "PERFORMANCE",
        f"{perf_score}/100",
        _grade_color(perf_score),
    )
    _pcard(
        slide,
        5.5,
        1.5,
        2.4,
        1.3,
        "SECURITY",
        f"{sec_score}/100",
        _grade_color(sec_score),
    )
    _pcard(
        slide,
        8.1,
        1.5,
        2.4,
        1.3,
        "ACCESSIBILITY",
        f"{a11y_score}/100",
        _grade_color(a11y_score),
    )
    _pcard(
        slide,
        10.7,
        1.5,
        2.4,
        1.3,
        "OVERALL",
        f"{fat_score}/100",
        _grade_color(fat_score),
    )

    summary = scores.get("summary", {})
    all_issues = []
    for prio_key in ["critical", "high", "medium", "low"]:
        for issue in summary.get(prio_key, []):
            all_issues.append(issue)

    _ptb(
        slide,
        0.5,
        3.2,
        12.3,
        0.4,
        f"Key Findings ({summary.get('issues_found', len(all_issues))} Issues)",
        20,
        P_GRAY,
        True,
    )
    y = 3.7
    prio_colors = (
        [P_RED] * len(summary.get("critical", []))
        + [P_ORANGE] * len(summary.get("high", []))
        + [P_YELLOW] * len(summary.get("medium", []))
        + [P_GREEN] * len(summary.get("low", []))
    )
    for i, issue in enumerate(all_issues[:10]):
        color = prio_colors[i] if i < len(prio_colors) else P_MGRAY
        _pbullet(slide, 0.7, y, issue[:100], 11, P_GRAY, color)
        y += 0.34

    # --- Chart slides (if charts exist) ---
    chart_slides = [
        ("chart_overview.png", "SEMrush Domain Overview"),
        ("chart_traffic_trend.png", "Organic Traffic Trend"),
        ("chart_keywords_trend.png", "Keyword Rankings Over Time"),
        ("chart_top_keywords.png", "Top Ranking Keywords"),
        ("chart_pagespeed.png", "PageSpeed Performance"),
        ("chart_fat_scores.png", "FAT Agent Score Summary"),
    ]

    for filename, title in chart_slides:
        if not charts_dir:
            continue
        chart_path = os.path.join(charts_dir, filename)
        if not os.path.exists(chart_path):
            continue

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide)
        _pheader(slide, title, brand)
        # Fit-to-box (aspect-ratio preserving) instead of a fixed W×H that
        # stretched charts whose native ratio differed from the slide.
        _add_chart_fit(slide, chart_path, 0.3, 1.3, 12.7, 5.2)
        source = f"Source: FAT Agent Audit | {url or 'Website'} | {date.today().strftime('%B %Y')}"
        _ptb(slide, 0.5, 6.7, 12.3, 0.4, source, 9, P_MGRAY, False, PP_ALIGN.CENTER)

    # --- SEO Priority Opportunities slide (SEMrush insight, not just numbers) ---
    if semrush and semrush.get("opportunity_keywords"):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide)
        _pheader(slide, "SEO Priority Opportunities", brand)
        _ptb(
            slide,
            0.5,
            1.15,
            12.3,
            0.35,
            "High-value keywords within striking distance — ranked by volume x CPC and how winnable the position is.",
            12,
            P_MGRAY,
            False,
        )
        rows = []
        for k in semrush["opportunity_keywords"][:12]:
            cpc = k.get("cpc")
            try:
                cpc_s = f"${float(cpc):.2f}"
            except (TypeError, ValueError):
                cpc_s = "-"
            rows.append(
                [
                    str(k.get("keyword", ""))[:40],
                    str(k.get("volume", "")),
                    cpc_s,
                    str(k.get("position", "")),
                    _strip_host(k.get("url", ""))[:30],
                    str(k.get("priority", "")),
                ]
            )
        _ptable(
            slide,
            0.4,
            1.65,
            12.5,
            ["Keyword", "Vol", "CPC", "Pos", "Target page", "Priority"],
            rows,
            col_widths=[3.3, 0.8, 0.9, 0.7, 3.4, 1.4],
            font_size=10,
            row_height=0.36,
        )

    # --- Keyword Cannibalisation slide ---
    if semrush and semrush.get("cannibalization"):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide)
        _pheader(slide, "Keyword Cannibalisation", brand)
        _ptb(
            slide,
            0.5,
            1.15,
            12.3,
            0.35,
            "Same keyword ranking on multiple URLs — the pages compete with each other and split ranking signals.",
            12,
            P_MGRAY,
            False,
        )
        rows = []
        for c in semrush["cannibalization"][:10]:
            urls = c.get("urls", []) or []
            count = len(urls) if urls else c.get("url_count", "")
            sample = ",  ".join(_strip_host(u) for u in urls[:3])
            rows.append(
                [
                    str(c.get("keyword", ""))[:38],
                    str(c.get("volume", "")),
                    str(count),
                    sample[:60],
                ]
            )
        _ptable(
            slide,
            0.4,
            1.65,
            12.5,
            ["Keyword", "Vol", "# URLs", "Competing pages (sample)"],
            rows,
            col_widths=[2.8, 0.8, 1.0, 5.6],
            font_size=10,
            row_height=0.42,
        )

    # --- Recommended Action Plan slide(s) — turns the report into next steps ---
    if recommendations:
        if isinstance(recommendations[0], dict):
            phases = [
                (r.get("phase") or r.get("title") or "", r.get("items") or [])
                for r in recommendations
            ]
        else:
            phases = [("Recommended Actions", list(recommendations))]
        phase_colors = [P_RED, P_ORANGE, P_GREEN, P_BLUE]

        def _new_action_slide():
            s = prs.slides.add_slide(prs.slide_layouts[6])
            add_bg(s)
            _pheader(s, "Recommended Action Plan", brand)
            return s

        slide = _new_action_slide()
        y = 1.4
        for pi, (pname, items) in enumerate(phases):
            color = phase_colors[pi % len(phase_colors)]
            if y > 6.4:
                slide = _new_action_slide()
                y = 1.4
            if pname:
                _ptb(slide, 0.5, y, 12.3, 0.4, pname, 16, color, True)
                y += 0.5
            for it in items:
                if y > 6.8:
                    slide = _new_action_slide()
                    y = 1.4
                _pbullet(slide, 0.7, y, str(it)[:115], 12, P_GRAY, color)
                y += 0.38
            y += 0.15

    # --- Closing slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, P_DARK)
    if brand and os.path.exists(brand):
        slide.shapes.add_picture(
            brand, PInches(4.9), PInches(0.5), PInches(3.5), PInches(3.5)
        )
    _ptb(slide, 0.5, 4.2, 12.3, 0.8, "Thank You", 40, P_WHITE, True, PP_ALIGN.CENTER)
    _ptb(
        slide,
        0.5,
        5.0,
        12.3,
        0.5,
        "FAT Agent Audit Complete",
        20,
        P_MGRAY,
        False,
        PP_ALIGN.CENTER,
    )
    _ptb(slide, 0.5, 5.6, 12.3, 0.5, url or "", 24, P_RED, True, PP_ALIGN.CENTER)
    _ptb(
        slide,
        0.5,
        6.3,
        12.3,
        0.5,
        f"Report Date: {date.today().strftime('%d %B %Y')}",
        14,
        P_MGRAY,
        False,
        PP_ALIGN.CENTER,
    )

    slug = safe_report_slug(url)
    filename = (
        f"FAT_Audit_Presentation_{slug}.pptx" if slug else "FAT_Audit_Presentation.pptx"
    )
    path = _safe_output_path(output_dir, filename)
    prs.save(path)
    return path


# ---------- Main ----------


def generate_html(
    scores, url, output_dir, charts_dir=None, client_facing=False, brand=None
):
    """Generate the self-contained HTML dashboard deliverable.

    Prefers importing generate_html_dashboard.generate_dashboard so the
    already-loaded (and optionally client-facing) scores dict is rendered
    in-process; falls back to subprocessing the script with --scores/--url/
    --output-dir/--charts-dir if the import is unavailable.
    """
    os.makedirs(output_dir, exist_ok=True)
    try:
        from generate_html_dashboard import generate_dashboard

        return generate_dashboard(
            scores=scores,
            url=url or "",
            output_dir=output_dir,
            charts_dir=charts_dir,
            client_facing=client_facing,
        )
    except ImportError:
        # Fallback: subprocess the dashboard script. It reads --scores from a
        # file, so persist the effective (possibly transformed) scores first.
        import tempfile

        script = os.path.join(
            os.path.dirname(os.path.abspath(__file__)), "generate_html_dashboard.py"
        )
        with tempfile.NamedTemporaryFile(
            "w", suffix=".json", delete=False, encoding="utf-8"
        ) as tf:
            json.dump(scores, tf)
            scores_path = tf.name
        try:
            cmd = [
                sys.executable,
                script,
                "--scores",
                scores_path,
                "--url",
                url or "",
                "--output-dir",
                output_dir,
            ]
            if charts_dir:
                cmd += ["--charts-dir", charts_dir]
            if client_facing:
                cmd.append("--client-facing")
            subprocess.run(cmd, check=True)
        finally:
            try:
                os.unlink(scores_path)
            except OSError:
                pass
        return os.path.join(output_dir, "FAT_Dashboard.html")


def main():
    if MISSING_DEPS:
        print(
            f"Error: Missing dependencies: {', '.join(MISSING_DEPS)}", file=sys.stderr
        )
        print(f"Install with: pip install {' '.join(MISSING_DEPS)}", file=sys.stderr)
        sys.exit(1)

    parser = argparse.ArgumentParser(description="Generate FAT Agent audit reports")
    parser.add_argument("--scores", help="Path to scored JSON file (default: stdin)")
    parser.add_argument("--semrush", help="Path to SEMrush data JSON file")
    parser.add_argument("--url", help="Audited site URL")
    parser.add_argument("--output-dir", default="./reports", help="Output directory")
    parser.add_argument(
        "--format", default="both", choices=["docx", "pptx", "html", "both"]
    )
    parser.add_argument("--charts-dir", help="Directory containing chart PNG images")
    parser.add_argument("--brand", help="Path to brand/logo image")
    parser.add_argument("--font", default="Plus Jakarta Sans", help="Font family")
    parser.add_argument("--pagespeed", help="Path to PageSpeed data JSON")
    parser.add_argument(
        "--actions",
        help="Path to a JSON action plan (list of strings, or list of "
        '{"phase","items":[...]} objects). Rendered as a Recommended Action '
        "Plan in the docx and pptx. Falls back to semrush['action_plan'] or "
        "scores['recommendations'] when omitted.",
    )
    parser.add_argument(
        "--client-facing",
        action="store_true",
        help="Rewrite findings/labels into plain-English, jargon-free, code-free language",
    )
    args = parser.parse_args()

    # Load scores
    scores = {}
    if args.scores:
        with open(args.scores, "r", encoding="utf-8") as f:
            scores = json.load(f)
    elif not sys.stdin.isatty():
        scores = json.load(sys.stdin)

    if not scores:
        print(
            "Error: No scores data provided. Pipe from calculate-score.py or use --scores",
            file=sys.stderr,
        )
        sys.exit(1)

    # Transform findings/labels for non-technical clients BEFORE rendering any
    # deliverable (docx, pptx, html).
    if args.client_facing:
        scores = make_client_facing(scores)

    semrush = None
    if args.semrush:
        with open(args.semrush, "r", encoding="utf-8") as f:
            semrush = json.load(f)

    pagespeed = None
    if args.pagespeed:
        with open(args.pagespeed, "r", encoding="utf-8") as f:
            pagespeed = json.load(f)

    # Action plan / recommendations: explicit --actions file wins, then an
    # action_plan embedded in the SEMrush JSON, then a recommendations key on
    # the scored JSON. This is what makes the deliverables actionable rather
    # than "just numbers".
    recommendations = None
    if args.actions:
        with open(args.actions, "r", encoding="utf-8") as f:
            recommendations = json.load(f)
    elif semrush and semrush.get("action_plan"):
        recommendations = semrush["action_plan"]
    elif isinstance(scores, dict) and scores.get("recommendations"):
        recommendations = scores["recommendations"]

    os.makedirs(args.output_dir, exist_ok=True)

    results = {}
    if args.format in ("docx", "both"):
        path = generate_docx(
            scores,
            args.url,
            args.output_dir,
            args.charts_dir,
            semrush,
            args.brand,
            pagespeed,
            args.font,
            recommendations=recommendations,
        )
        results["docx"] = path
        print(f"Word report: {path}")

    if args.format in ("pptx", "both"):
        path = generate_pptx(
            scores,
            args.url,
            args.output_dir,
            args.charts_dir,
            semrush,
            args.brand,
            pagespeed,
            args.font,
            recommendations=recommendations,
        )
        results["pptx"] = path
        print(f"PowerPoint:  {path}")

    if args.format == "html":
        path = generate_html(
            scores,
            args.url,
            args.output_dir,
            charts_dir=args.charts_dir,
            client_facing=args.client_facing,
        )
        results["html"] = path
        print(f"HTML dashboard: {path}")

    print(json.dumps(results, indent=2))


if __name__ == "__main__":
    main()
